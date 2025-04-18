import io
import os

import pdfplumber
import pytesseract
from docx import Document
from PIL import Image
from pptx import Presentation

# Set the path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    text = ""
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise ValueError(f"PDF text extraction failed: {str(e)}")


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise ValueError(f"DOCX text extraction failed: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"PPTX text extraction failed: {str(e)}")


def extract_text_from_image(image_path):
    """Extract text from an image file using Tesseract OCR"""
    try:
        # Check if the image file exists
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"Image file not found: {image_path}")

        # Extract text from the image
        text = pytesseract.image_to_string(image_path)
        return text.strip()
    except Exception as e:
        raise Exception(f"Error extracting text from image: {str(e)}")


def extract_text_from_file(file_path: str) -> dict:
    """Main function to extract text from various file types"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return {"text": extract_text_from_pdf(file_path), "file_type": "pdf"}
    elif ext == ".docx":
        return {"text": extract_text_from_docx(file_path), "file_type": "docx"}
    elif ext == ".pptx":
        return {"text": extract_text_from_pptx(file_path), "file_type": "pptx"}
    elif ext in (".jpg", ".jpeg", ".png", ".bmp"):
        return {"text": extract_text_from_image(file_path), "file_type": "image"}
    else:
        raise ValueError(f"Unsupported file type: {ext}")
